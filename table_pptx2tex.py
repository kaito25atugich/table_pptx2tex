from pptx import Presentation
import re

def pptx_table_to_latex(table):
    latex_code = "\\begin{table}[ht]\n"
    latex_code += "\\centering\n"
    latex_code += "\\begin{tabular}{|" + "|".join(["c"] * len(table.columns)) + "|}\n"
    latex_code += "\\hline\n"

    # ヘッダー行の処理
    latex_code += " & ".join(table.cell(0, col_idx).text for col_idx in range(len(table.columns))) + " \\\\\n"
    latex_code += "\\hline\n"

    # データ行の処理
    for row in table.rows:
        latex_code += " & ".join(cell.text for cell in row.cells) + " \\\\\n"

    latex_code += "\\hline\n"
    latex_code += "\\end{tabular}\n"
    latex_code += "\\caption{Your table caption here.}\n"
    latex_code += "\\label{tab:your_table_label}\n"
    latex_code += "\\end{table}\n"

    return latex_code

def pptx_to_latex(pptx_file, output_path):
    prs = Presentation(pptx_file)

    with open(output_file, 'w') as f:
        for slide_number, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_table:
                    # テーブルを見つけた場合
                    table = shape.table
                    latex_code = pptx_table_to_latex(table)
                    shape._element.getparent().remove(shape._element)
                    f.write(latex_code + '\n')

if __name__ == "__main__":
    pptx_file = "1218.pptx"  # ご自身のPowerPointファイル名に変更してください
    output_file = "output_table.txt"
    pptx_to_latex(pptx_file, output_file)